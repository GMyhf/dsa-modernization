"""Courseware gate regressions; run with python -m unittest discover -s courseware."""
import tempfile
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

import verify
from pptx import Presentation
from pptx.dml.color import RGBColor


class GateRegressionTests(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.addCleanup(self.tmp.cleanup)
        self.root = Path(self.tmp.name)
        self.stem = verify.build_all.CHAPTERS['01']
        self.pptx = self.root / (self.stem + '.pptx')
        self.meta = dict(title='Test', subtitle='Subtitle', footer='Footer', info='Info')
        self.modules = {'01': SimpleNamespace(META=self.meta, SLIDES=[])}
        verify.deck.build(self.meta, [], str(self.pptx))
        (self.root / 'README.md').write_text(
            f'| 1 | `{self.stem}` | 1 |\n', encoding='utf-8')
        self.patch = patch.object(verify, 'HERE', self.root)
        self.patch.start()
        self.addCleanup(self.patch.stop)
        verify.FAILURES.clear()
        verify.NOTES.clear()

    def test_unchanged_deck_passes(self):
        verify.check_regenerate(['01'], self.modules)
        self.assertEqual(verify.FAILURES, [])

    def test_moved_shape_is_rejected(self):
        prs = Presentation(self.pptx)
        prs.slides[0].shapes[0].left += 914400
        prs.save(self.pptx)
        verify.check_regenerate(['01'], self.modules)
        self.assertTrue(verify.FAILURES, 'Changed layout must require regeneration')

    def test_changed_background_is_rejected(self):
        prs = Presentation(self.pptx)
        fill = prs.slides[0].background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 0, 0)
        prs.save(self.pptx)
        verify.check_regenerate(['01'], self.modules)
        self.assertTrue(verify.FAILURES)

    def test_zip_timestamps_do_not_require_regeneration(self):
        contents = verify.package_contents(self.pptx)
        with zipfile.ZipFile(self.pptx, 'w') as package:
            for name, data in contents.items():
                info = zipfile.ZipInfo(name, date_time=(2000, 1, 1, 0, 0, 0))
                package.writestr(info, data)
        verify.check_regenerate(['01'], self.modules)
        self.assertEqual(verify.FAILURES, [])

    def render(self, bbox, returncode=0):
        def run(cmd, **kwargs):
            output = ''
            code = 0
            if '--convert-to' in cmd:
                out = Path(cmd[cmd.index('--outdir') + 1])
                (out / (self.stem + '.pdf')).touch()
            elif cmd[0] == 'pdffonts':
                output = 'name type encoding emb sub uni object ID\n'
                output += 'ABC+NotoSansCJKsc Type 1 Identity-H yes yes yes 1 0\n'
            elif '-bbox' in cmd:
                output, code = bbox, returncode
            return SimpleNamespace(stdout=output, stderr='probe error', returncode=code)
        with patch.object(verify.subprocess, 'run', side_effect=run):
            verify.check_render(['01'])

    def test_text_extraction_failure_is_rejected(self):
        self.render('', returncode=1)
        self.assertTrue(verify.FAILURES, 'Failed extraction must not pass as zero pages')

    def test_empty_extraction_is_rejected(self):
        self.render('')
        self.assertTrue(verify.FAILURES)

    def test_extra_page_is_rejected(self):
        self.render('<page width="960" height="540"></page>' * 2)
        self.assertTrue(verify.FAILURES)

    def test_valid_render_passes(self):
        self.render('<page width="960" height="540">'
                    '<word xMin="20" yMin="20" xMax="40" yMax="40">Test</word>'
                    '</page>')
        self.assertEqual(verify.FAILURES, [])


if __name__ == '__main__':
    unittest.main()
