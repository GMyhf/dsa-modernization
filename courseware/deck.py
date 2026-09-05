# -*- coding: utf-8 -*-
"""构建课件 PPTX 的小工具库。

设计目标：
  * 16:9 版面，中文字体（微软雅黑）在 Windows / Office 下开箱即用；
  * 用一组简单的"幻灯片描述"元组来写内容，与排版代码解耦；
  * 代码页按行数自动缩放字号，尽量避免溢出。

幻灯片描述（见各 content/wNN.py）支持的类型：
  ("title",   主标题, 副标题)
  ("section", 编号, 章节标题, 可选小字)
  ("bullets", 标题, [条目, ...])            条目以 "- " 前缀表示次级
  ("code",    标题, 代码字符串, 说明)
  ("table",   标题, [[表头...], [行...], ...], 说明)
  ("two",     标题, 左标题, [左条目...], 右标题, [右条目...])
  ("key",     标题, 要点正文)               整页强调一句话
  ("ascii",   标题, 等宽示意图, 说明)
  ("image",   标题, 图片路径, 说明)      原书扫描图，等比缩放后居中
"""

import math
import pathlib
import re

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 主题配置
EA_FONT = '微软雅黑'          # 中文字体
LATIN_FONT = 'Segoe UI'       # 西文字体
MONO_FONT = 'Consolas'        # 代码字体

NAVY = RGBColor(0x12, 0x39, 0x5B)      # 主色：深蓝
ACCENT = RGBColor(0xE0, 0x7B, 0x39)    # 强调色：橙
INK = RGBColor(0x24, 0x2A, 0x33)       # 正文
MUTED = RGBColor(0x6B, 0x75, 0x82)     # 次要文字
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)     # 浅底
CODE_BG = RGBColor(0xF6, 0xF7, 0xF9)
RULE = RGBColor(0xD6, 0xDC, 0xE4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
BODY_TOP = Inches(1.42)
BODY_H = Inches(5.32)
BODY_W = SLIDE_W - 2 * MARGIN


# ---------------------------------------------------------------- 字体工具
def _style_run(run, size, bold=False, color=INK, mono=False, italic=False):
    """设置一个 run 的字体。python-pptx 只写 a:latin，中文需手工补 a:ea。"""
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = MONO_FONT if mono else LATIN_FONT
    rPr = f._rPr
    for tag in ('a:ea', 'a:cs'):
        for old in rPr.findall(qn(tag)):
            rPr.remove(old)
    latin = rPr.find(qn('a:latin'))
    ea = rPr.makeelement(qn('a:ea'), {'typeface': MONO_FONT if mono else EA_FONT})
    if latin is not None:
        latin.addnext(ea)
    else:
        rPr.append(ea)


def _textbox(slide, left, top, width, height, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    # 固定尺寸：否则不换行的文本框会被渲染器自动居中，左边距忽宽忽窄
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return box, tf


# 强调段内部允许**单个** `*`（`[[0]*n]*m` 这类代码片段里全是它），
# 但不允许跨过下一个 `**` —— 否则会把两段强调粘成一段。
# 旧写法 `\*\*[^*]+\*\*` 匹配不上含 `*` 的内容，于是 `**` 自己被原样印了出来。
_SEGMENT = re.compile(r'(\*\*(?:[^*]|\*(?!\*))+?\*\*|`[^`]+`)')


def _plain(text):
    """去掉 **强调** 与 `等宽` 标记，用于度量文字宽度。"""
    return text.replace('**', '').replace('`', '')


_INLINE_CODE = re.compile(r'(`[^`]+`)')


def _add_runs(p, text, size, bold, color, mono):
    r"""把一行文本按 **强调** / `等宽` 切成多个 run。

    代码/示意图（mono）原样输出：Python 的 ** 幂运算符不能被当成强调标记。

    ⚠️ **强调段内部必须继续拆 `等宽`**。`_SEGMENT` 的 `\*\*[^*]+\*\*`
    会把 ``**负数不能用 `bin(x)`**`` 整段吞掉；若在这里直接 `seg[2:-2]` 输出，
    反引号就原样印在放映稿上了。cs101 曾有 8 处这种写法，
    在 T-007（439 页）、T-011（8 页）、T-013（29 页）三轮逐页复核里**全部漏过** ——
    人眼对一个小小的反引号不敏感，所以这件事必须由代码保证。
    """
    if mono:
        run = p.add_run()
        run.text = text
        _style_run(run, size, bold, color, True)
        return

    def emit(chunk, chunk_bold, chunk_color):
        for part in _INLINE_CODE.split(chunk):
            if not part:
                continue
            if part.startswith('`') and part.endswith('`') and len(part) > 1:
                run = p.add_run()
                run.text = part[1:-1]
                _style_run(run, size * 0.94, chunk_bold, chunk_color, True)
            else:
                run = p.add_run()
                run.text = part
                _style_run(run, size, chunk_bold, chunk_color, mono)

    for seg in _SEGMENT.split(text):
        if not seg:
            continue
        if seg.startswith('**') and seg.endswith('**'):
            emit(seg[2:-2], True, NAVY if color is INK else color)
        else:
            emit(seg, bold, color)


def _para(tf, text, size, bold=False, color=INK, mono=False,
          space_before=0, space_after=6, first=False, align=PP_ALIGN.LEFT,
          line_spacing=None, indent_level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    p.level = indent_level
    _add_runs(p, text, size, bold, color, mono)
    return p


def _rect(slide, left, top, width, height, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


# ---------------------------------------------------------------- 排版度量
def _ems(text):
    """估算一段文字的宽度，单位为 em（1 em = 当前字号）。

    CJK 与全角标点按 1.0 计，ASCII 按 0.55 计 —— 用于估算换行与列宽。
    """
    w = 0.0
    for ch in _plain(text):
        w += 1.0 if ord(ch) > 0x2E80 else 0.55
    return w


def _wrapped_lines(text, size_pt, width_pt):
    """给定字号与可用宽度，估算这段文字会占几行。"""
    if not text:
        return 1
    return max(1, math.ceil(_ems(text) * size_pt / width_pt))


def _fit_size(items, width_pt, height_pt, hi, lo, line_spacing, gap_ratio):
    """在 [lo, hi] 中挑最大的字号，使全部条目仍能装进 height_pt。

    items 为 (文本, 相对字号比例) 的列表。
    """
    size = hi
    while size > lo:
        total = 0.0
        for text, ratio in items:
            s = size * ratio
            total += _wrapped_lines(text, s, width_pt) * s * line_spacing
            total += s * gap_ratio
        if total <= height_pt:
            break
        size -= 0.5
    return size



# ---------------------------------------------------------------- 溢出保护
MIN_PT = 8.5          # 代码 / 示意图页的字号下限


class LayoutOverflow(ValueError):
    """内容多到自适应字号已触底 —— 再缩就不可读，再画就会溢出版心。"""


def _guard_min_size(kind, title, nlines, widest, want, floor):
    """字号触底时**直接报错**，而不是静默钳位后照常画。

    历史教训：这里原先只写 `size = max(8.5, ...)`。字号被钳到下限后文字照常输出，
    盒子却被 clamp 在版心内 —— 于是代码压到页脚上，而
      * 前 7 项检查（配对 / 元数据 / 语法 / 可重生成 …）根本看不到版面；
      * 渲染检查当时只比对**页面**边界，压到页脚仍在页内，也报"0 处越界"。
    结果 W16 有 3 页长期是坏的。字号是构建期就能精确算出来的量，
    应当在这里当场失败，而不是留给 4 分钟的渲染检查去碰运气。
    """
    if want >= floor:
        return
    raise LayoutOverflow(
        f'{kind} 页「{title}」内容过多：{nlines} 行、最长 {widest:.0f} em，'
        f'需要 {want:.2f}pt 才能装进版心，已低于下限 {floor}pt。\n'
        f'        请拆成两页或精简内容 —— 不要调低下限。'
    )


# ---------------------------------------------------------------- 版面构件
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _slide_header(slide, title):
    """页眉：标题 + 橙色短线。"""
    _, tf = _textbox(slide, MARGIN, Inches(0.46), BODY_W, Inches(0.6))
    _para(tf, title, 26, bold=True, color=NAVY, first=True, space_after=0)
    _rect(slide, MARGIN, Inches(1.16), Inches(1.05), Pt(3.2), fill=ACCENT)


def _footer(slide, label, number):
    _, tf = _textbox(slide, MARGIN, Inches(6.92), Inches(9.0), Inches(0.34))
    _para(tf, label, 10, color=MUTED, first=True, space_after=0)
    _, tf2 = _textbox(slide, SLIDE_W - MARGIN - Inches(1.0), Inches(6.92),
                      Inches(1.0), Inches(0.34))
    _para(tf2, str(number), 10, color=MUTED, first=True, space_after=0,
          align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- 各类页面
def _add_title(prs, main, sub, meta):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    _rect(slide, 0, Inches(4.05), SLIDE_W, Pt(3.2), fill=ACCENT)
    _, tf = _textbox(slide, MARGIN, Inches(2.30), SLIDE_W - 2 * MARGIN, Inches(1.7))
    _para(tf, main, 40, bold=True, color=WHITE, first=True, space_after=10)
    if sub:
        _para(tf, sub, 19, color=RGBColor(0xC7, 0xD3, 0xE2), space_after=0)
    _, tf2 = _textbox(slide, MARGIN, Inches(4.55), SLIDE_W - 2 * MARGIN, Inches(1.6))
    firstline = True
    for line in meta:
        _para(tf2, line, 13, color=RGBColor(0xA8, 0xB8, 0xCC),
              first=firstline, space_after=5)
        firstline = False
    return slide


def _add_section(prs, number, title, note=''):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT)
    _rect(slide, 0, Inches(2.55), Inches(0.30), Inches(1.9), fill=ACCENT)
    _, tf = _textbox(slide, Inches(0.95), Inches(2.62), Inches(11.4), Inches(1.9))
    _para(tf, number, 15, bold=True, color=ACCENT, first=True, space_after=8)
    _para(tf, title, 34, bold=True, color=NAVY, space_after=8)
    if note:
        _para(tf, note, 14, color=MUTED, space_after=0)
    return slide


def _add_bullets(prs, title, items):
    slide = _blank(prs)
    _slide_header(slide, title)
    body_w_pt = BODY_W / 12700.0
    body_h_pt = BODY_H / 12700.0

    parsed = []
    for raw in items:
        sub = raw.startswith('- ')
        text = raw[2:] if sub else raw
        head = (text.startswith('**') and text.endswith('**')
                and text.count('**') == 2)
        if head:
            text = text[2:-2]
        parsed.append((text, sub, head))

    size = _fit_size([(t, 0.82 if sub else 1.0) for t, sub, _ in parsed],
                     body_w_pt, body_h_pt * 0.92, hi=20, lo=11.5,
                     line_spacing=1.32, gap_ratio=0.62)
    gap = size * 0.62

    _, tf = _textbox(slide, MARGIN, BODY_TOP, BODY_W, BODY_H)
    first = True
    for text, sub, bold in parsed:
        if sub:
            _para(tf, '·  ' + text, size * 0.82, color=MUTED, first=first,
                  space_after=gap * 0.72, indent_level=1, line_spacing=1.3)
        else:
            _para(tf, ('▍ ' if bold else '•  ') + text, size, bold=bold,
                  color=NAVY if bold else INK, first=first,
                  space_after=gap, line_spacing=1.32)
        first = False
    return slide


def _add_code(prs, title, code, caption=''):
    slide = _blank(prs)
    _slide_header(slide, title)
    lines = code.rstrip('\n').split('\n')
    n = len(lines)
    widest = max((_ems(l) for l in lines), default=1)

    cap_h = Inches(0.46) if caption else Inches(0)
    avail_h_pt = (BODY_H - cap_h) / 12700.0
    avail_w_pt = (BODY_W - Inches(0.72)) / 12700.0

    # 行距用绝对磅值，保证盒子高度可精确计算，不会溢出
    LEAD = 1.30                                   # 行高 / 字号
    PAD = 22.0                                    # 盒子上下内边距合计（pt）
    by_rows = (avail_h_pt - PAD) / (n * LEAD)
    by_cols = avail_w_pt / max(widest, 1) / 0.60  # 等宽字体约 0.60 em/字符
    _guard_min_size('code', title, n, widest, min(by_rows, by_cols), MIN_PT)
    size = max(MIN_PT, min(15.5, by_rows, by_cols))

    box_h_pt = min(float(BODY_H) / 12700.0 - float(cap_h) / 12700.0,
                   n * size * LEAD + PAD)
    box_h = Emu(int(box_h_pt * 12700))
    _rect(slide, MARGIN, BODY_TOP, BODY_W, box_h, fill=CODE_BG, line=RULE)
    _, tf = _textbox(slide, MARGIN + Inches(0.28), BODY_TOP + Inches(0.11),
                     BODY_W - Inches(0.56), box_h - Inches(0.22))
    first = True
    for line in lines:
        _para(tf, line if line else ' ', size, mono=True, color=INK,
              first=first, space_after=0, line_spacing=Pt(size * LEAD))
        first = False
    if caption:
        _, tf2 = _textbox(slide, MARGIN, BODY_TOP + box_h + Inches(0.13),
                          BODY_W, Inches(0.36))
        _para(tf2, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


def _add_ascii(prs, title, art, caption=''):
    slide = _blank(prs)
    _slide_header(slide, title)
    lines = art.rstrip('\n').split('\n')
    n = len(lines)
    widest = max((_ems(l) for l in lines), default=1)
    cap_h = Inches(0.46) if caption else Inches(0)
    avail_h_pt = (BODY_H - cap_h) / 12700.0
    avail_w_pt = BODY_W / 12700.0

    LEAD = 1.30
    by_rows = avail_h_pt / (n * LEAD)
    by_cols = avail_w_pt / max(widest, 1) / 0.60
    _guard_min_size('ascii', title, n, widest, min(by_rows, by_cols), MIN_PT)
    size = max(MIN_PT, min(17.0, by_rows, by_cols))

    art_h_pt = n * size * LEAD
    top = BODY_TOP + Emu(int(max(0.0, (avail_h_pt - art_h_pt) / 2) * 12700))
    _, tf = _textbox(slide, MARGIN, top, BODY_W,
                     Emu(int(art_h_pt * 12700)))
    first = True
    for line in lines:
        _para(tf, line if line else ' ', size, mono=True, color=NAVY,
              first=first, space_after=0, line_spacing=Pt(size * LEAD))
        first = False
    if caption:
        _, tf2 = _textbox(slide, MARGIN, Inches(6.44), BODY_W, Inches(0.36))
        _para(tf2, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


def _add_image(prs, title, path, caption=''):
    """整页放一张原书插图：等比缩放到版心，居中。

    图注与 ascii/code 页一样固定在版心底部——它常常是"这张图要看什么"的
    唯一说明，不能因为图高不同而上下乱跑。
    """
    slide = _blank(prs)
    _slide_header(slide, title)
    src = pathlib.Path(path)
    if not src.is_absolute():
        src = (pathlib.Path(__file__).resolve().parent / src)
    if not src.exists():
        raise FileNotFoundError(f'image 页「{title}」找不到图片：{path}')

    from PIL import Image  # 仅为读出像素尺寸；缺库时退回按 4:3 估算
    try:
        px_w, px_h = Image.open(src).size
    except Exception:                                    # pragma: no cover
        px_w, px_h = 4, 3

    cap_h = Inches(0.52) if caption else Inches(0)
    avail_w = BODY_W
    avail_h = BODY_H - cap_h - Inches(0.10)
    scale = min(float(avail_w) / px_w, float(avail_h) / px_h)
    w, h = Emu(int(px_w * scale)), Emu(int(px_h * scale))
    left = MARGIN + Emu(int((float(avail_w) - float(w)) / 2))
    top = BODY_TOP + Emu(int((float(avail_h) - float(h)) / 2))
    # 扫描件是白底黑线，铺一层白底 + 细边框，投影时不会和页面浅灰糊在一起
    _rect(slide, left - Inches(0.10), top - Inches(0.10),
          w + Inches(0.20), h + Inches(0.20), fill=WHITE, line=RULE)
    slide.shapes.add_picture(str(src), left, top, width=w, height=h)
    if caption:
        _, tf = _textbox(slide, MARGIN, Inches(6.44), BODY_W, Inches(0.36))
        _para(tf, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


# 表格行高估算的两个经验系数。它们都来自一次实际的排版事故，不是保守起见随手加的。
#
# ⚠️ 事故：第 1 章「六种算法设计方法」那一页，图注**压在表格最后一行的字上**。
#    越界检查看不到这种事 —— 文字仍在版心内，只是压在别的文字上。两处都算少了：
#
#  * TABLE_SAFETY —— 「动态规划」那一格 4 个汉字、列宽 83pt、字号 17pt，
#    算出来 4 x 17 = 68pt，可用 67pt 上下，**富余不到 2pt**，估算判「一行装得下」；
#    渲染器实际换成了两行。汉字宽度按 1.0 em 估已经接近真值，差的就是这点余量。
#  * TABLE_LEAD —— 行高原本按「字号 x 1.30」算。中文字体的自然行高远大于字号
#    （Noto Sans CJK 约 1.48 em），乘上 1.22 的行距后一行要 1.5 em 上下，
#    于是每个两行的单元格都比预留的高 13%，六行累积下来就是小半厘米。
#
# 这两个数都偏保守：宁可字号小半档，也不要让图注压到表格上。
TABLE_SAFETY = 1.12      # 估算换行时给列宽留的余量
TABLE_LEAD = 1.55        # 行高 / 字号（含中文字体的自然行高与行距）


def _row_heights(rows, col_w, size, pad=16.0):
    """给定字号，逐行估算行高（pt）。"""
    heights = []
    for row in rows:
        lines = max(
            _wrapped_lines(str(row[c]), size * TABLE_SAFETY, col_w[c] - pad)
            for c in range(len(row)))
        heights.append(max(lines * size * TABLE_LEAD + 12, 34.6))
    return heights


def _add_table(prs, title, rows, caption=''):
    slide = _blank(prs)
    _slide_header(slide, title)
    nrow, ncol = len(rows), len(rows[0])
    cap_h = Inches(0.62) if caption else Inches(0)
    avail_h_pt = (BODY_H - cap_h) / 12700.0
    total_w_pt = BODY_W / 12700.0

    # 列宽按各列最长内容分配，并做上下限收敛，避免某一列被挤到换行
    weights = []
    for c in range(ncol):
        w = max(_ems(str(rows[r][c])) for r in range(nrow))
        weights.append(min(max(w, 4.0), 34.0))
    tot = sum(weights)
    col_w = [total_w_pt * w / tot for w in weights]

    # 在能装下的前提下取最大字号
    size, floor = 17.0, 9.0
    while size > floor:
        if sum(_row_heights(rows, col_w, size)) <= avail_h_pt:
            break
        size -= 0.5
    else:
        raise LayoutOverflow(
            f'table 页「{title}」内容过多：{nrow} 行 x {ncol} 列，'
            f'字号降到下限 {floor}pt 仍装不进版心。\n'
            f'        请拆成两页或精简单元格 —— 不要调低下限。')

    row_h = _row_heights(rows, col_w, size)
    height = Emu(int(sum(row_h) * 12700))

    shape = slide.shapes.add_table(nrow, ncol, MARGIN, BODY_TOP, BODY_W, height)
    table = shape.table
    for c in range(ncol):
        table.columns[c].width = Emu(int(col_w[c] * 12700))
    for r, row in enumerate(rows):
        table.rows[r].height = Emu(int(row_h[r] * 12700))
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.11)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (
                WHITE if r % 2 else LIGHT)
            tf = cell.text_frame
            tf.word_wrap = True
            _para(tf, str(val), size, bold=(r == 0),
                  color=WHITE if r == 0 else INK, first=True, space_after=0,
                  line_spacing=1.22)
    if caption:
        # 固定在版心底部：估算再准也只是估算，渲染器的行高各不相同，
        # 所以图注不跟着表格走。上面的 TABLE_SAFETY / TABLE_LEAD 负责让
        # 表格**够不到**这一带；verify.py 第 7 项负责在渲染后复核这件事。
        top = min(BODY_TOP + height + Inches(0.14), Inches(6.38))
        _, tf2 = _textbox(slide, MARGIN, top, BODY_W, Inches(0.36))
        _para(tf2, caption, 13, color=MUTED, first=True, space_after=0)
    return slide


def _add_two(prs, title, lhead, litems, rhead, ritems):
    slide = _blank(prs)
    _slide_header(slide, title)
    colw = (BODY_W - Inches(0.5)) / 2
    for i, (head, items) in enumerate(((lhead, litems), (rhead, ritems))):
        left = MARGIN + Emu(int(i * (colw + Inches(0.5))))
        _rect(slide, left, BODY_TOP, Emu(int(colw)), Inches(0.46), fill=NAVY)
        _, th = _textbox(slide, left + Inches(0.16), BODY_TOP + Inches(0.09),
                         Emu(int(colw)) - Inches(0.32), Inches(0.34))
        _para(th, head, 16, bold=True, color=WHITE, first=True, space_after=0)
        n = len(items)
        size = 16 if n <= 5 else (14.5 if n <= 7 else 13)
        _, tf = _textbox(slide, left + Inches(0.10), BODY_TOP + Inches(0.66),
                         Emu(int(colw)) - Inches(0.20), BODY_H - Inches(0.8))
        first = True
        for raw in items:
            # 整条都用反引号包住 = 整行等宽；否则原样交给 _add_runs，
            # 由它去拆行内的 **强调** 与 `等宽`。
            # ⚠️ 旧写法是 `raw.strip('`')` —— 对 "`[[nodiscard]]`：返回值丢不掉"
            # 这种「行内代码开头、后面还有中文」的条目，它剥掉开头那个反引号、
            # 留下中间那个，于是反引号**原样印在放映稿上**。
            mono = len(raw) > 1 and raw.startswith('`') and raw.endswith('`')
            text = raw[1:-1] if mono else raw
            _para(tf, ('· ' if not mono else '') + text, size, mono=mono,
                  color=INK, first=first, space_after=8, line_spacing=1.25)
            first = False
    return slide


def _add_key(prs, title, text):
    slide = _blank(prs)
    _slide_header(slide, title)
    _rect(slide, MARGIN, Inches(2.28), BODY_W, Inches(2.5), fill=LIGHT)
    _rect(slide, MARGIN, Inches(2.28), Inches(0.09), Inches(2.5), fill=ACCENT)
    _, tf = _textbox(slide, MARGIN + Inches(0.55), Inches(2.60),
                     BODY_W - Inches(1.1), Inches(1.9))
    size = 26 if len(text) <= 40 else (22 if len(text) <= 70 else 18)
    _para(tf, text, size, bold=True, color=NAVY, first=True, space_after=0,
          line_spacing=1.45)
    return slide


# ---------------------------------------------------------------- 构建入口
_BUILDERS = {
    'section': lambda prs, s: _add_section(prs, s[1], s[2],
                                           s[3] if len(s) > 3 else ''),
    'bullets': lambda prs, s: _add_bullets(prs, s[1], s[2]),
    'code': lambda prs, s: _add_code(prs, s[1], s[2],
                                     s[3] if len(s) > 3 else ''),
    'ascii': lambda prs, s: _add_ascii(prs, s[1], s[2],
                                       s[3] if len(s) > 3 else ''),
    'image': lambda prs, s: _add_image(prs, s[1], s[2],
                                       s[3] if len(s) > 3 else ''),
    'table': lambda prs, s: _add_table(prs, s[1], s[2],
                                       s[3] if len(s) > 3 else ''),
    'two': lambda prs, s: _add_two(prs, s[1], s[2], s[3], s[4], s[5]),
    'key': lambda prs, s: _add_key(prs, s[1], s[2]),
}


def build(meta, slides, out_path):
    """meta = {'title','subtitle','footer','info':[...]}；slides 见模块文档。"""
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _add_title(prs, meta['title'], meta.get('subtitle', ''), meta.get('info', []))
    for i, spec in enumerate(slides, start=2):
        kind = spec[0]
        if kind == 'title':
            continue
        slide = _BUILDERS[kind](prs, spec)
        if kind != 'section':
            _footer(slide, meta.get('footer', ''), i)
    prs.save(out_path)
    return len(prs.slides)
