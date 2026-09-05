#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""把课件 + 旁白讲稿合成为带解说的 mp4。

用法:
    python3 make_video.py 01                # 合成第 1 章
    python3 make_video.py 01 --audio-only   # 只跑 TTS，看看总时长够不够一节课
    python3 make_video.py 01 --check        # 只校验产物是否最新（闸门用）
    python3 make_video.py 01 --preview      # 从成品转一份 720p 轻量版，便于传阅

画面**从当前 .pptx 现导**（LibreOffice → PDF → PNG），不依赖任何预先存在的图片，
所以视频永远和课件同版。旁白来自 content/chNN_narration.md，一节对应一页。

产物:
    DSA_CHNN_*.mp4                 成品（1080p）
    video/chNN-preview.mp4         轻量版（720p 单声道，--preview 才生成）
    video/chNN.srt                 字幕（每页一条）
    video/chNN.timeline.json       时间轴 + 两个来源文件的 sha256
    content/chNN_narration.md      文末「时间控制表」按实测时长重写

环境变量:
    TTS_VOICE   默认 zh-CN-YunyangNeural
    TTS_RATE    默认 +0%
    VIDEO_WORK  中间产物目录，默认 courseware/.videowork（已 gitignore）

⚠️ 语音按**讲稿文本的哈希**缓存，不是「文件在就复用」。
   讲稿改了一页，只有那一页重新合成；而如果只按文件是否存在判断，
   改了词却照旧用旧音频 —— 视频和讲稿从此对不上，且没有任何报错。
"""

import argparse
import hashlib
import io
import json
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))
import build_all                                            # noqa: E402

WORK = Path(os.environ.get('VIDEO_WORK', HERE / '.videowork'))
VOICE = os.environ.get('TTS_VOICE', 'zh-CN-YunyangNeural')
RATE = os.environ.get('TTS_RATE', '+0%')

HEAD = 0.4      # 每页开始前的静默
TAIL = 1.0      # 每页读完后的静默留白
BG = '0x12395B'  # 画面补边色，取 deck.py 的主色 NAVY

SECTION = re.compile(r'^## P(\d+)｜(.+?)\s*$', re.M)


# ---------------------------------------------------------------- 小工具
def run(cmd, **kw):
    r = subprocess.run([str(c) for c in cmd], capture_output=True, **kw)
    if r.returncode != 0:
        sys.exit('FAILED: %s\n%s' % (' '.join(map(str, cmd)),
                                     r.stderr.decode('utf-8', 'replace')[-2000:]))
    return r.stdout.decode('utf-8', 'replace')


def pptx_digest(path):
    """课件**画面内容**的摘要 —— 不是文件字节的摘要。

    ⚠️ python-pptx 写出来的 .pptx **不是逐字节可复现的**（zip 里带时间戳）：
    同一份 content 连着生成两次，字节摘要就不一样。用它当「视频过没过期」的
    判据，等于每次重建课件都谎报一次过期 —— 判据一旦学会撒谎，就没人再看它了。
    （这和讲稿时间控制表那次是同一类错误：把一个不稳定的量当成了指纹。）

    这里改成摘要**决定画面的那些东西**：每页每个形状的位置尺寸、文字、
    每个 run 的字号/加粗/字体，表格逐格文字，以及图片的内容哈希。
    渲染出来长得一样，摘要就一样。
    """
    from pptx import Presentation
    h = hashlib.sha256()
    for slide in Presentation(str(path)).slides:
        h.update(b'\x00slide')
        for shape in slide.shapes:
            h.update(('|%s,%s,%s,%s,%s' % (shape.shape_type, shape.left, shape.top,
                                           shape.width, shape.height)).encode())
            if shape.shape_type == 13:            # PICTURE
                h.update(shape.image.blob)
            frames = []
            if shape.has_text_frame:
                frames.append(shape.text_frame)
            elif shape.has_table:
                frames += [c.text_frame for r in shape.table.rows for c in r.cells]
            for tf in frames:
                for para in tf.paragraphs:
                    for run in para.runs:
                        f = run.font
                        h.update(('|%s;%s;%s;%s' % (
                            run.text, f.size, f.bold, f.name)).encode())
    return h.hexdigest()


def duration(path):
    return float(run(['ffprobe', '-v', 'error', '-show_entries', 'format=duration',
                      '-of', 'default=nw=1:nk=1', path]).strip())


def fmt(seconds):
    s = int(round(seconds))
    return '%02d:%02d' % (s // 60, s % 60)


def srt_ts(seconds):
    ms = int(round(seconds * 1000))
    return '%02d:%02d:%02d,%03d' % (ms // 3600000, ms // 60000 % 60,
                                    ms // 1000 % 60, ms % 1000)


# ---------------------------------------------------------------- 讲稿
def parse_narration(path, expected_pages):
    """把讲稿切成每页一段朗读文本。"""
    md = io.open(path, encoding='utf-8').read()
    md = md.split('## 附：时间控制表')[0]
    parts = SECTION.split(md)
    sections = []
    for i in range(1, len(parts), 3):
        page, title, body = int(parts[i]), parts[i + 1].strip(), parts[i + 2]
        lines = []
        for ln in body.split('\n'):
            ln = ln.strip()
            # 跳过标题、引用、表格、列表和分隔线 —— 那些是给读者看的，不朗读
            if not ln or ln.startswith(('#', '>', '|', '-', '*', '<!--')):
                continue
            ln = re.sub(r'【[^】]*】', '', ln)      # 翻页提示不朗读
            ln = ln.replace('**', '')
            ln = re.sub(r'[ \t　]+', ' ', ln).strip()
            if ln:
                lines.append(ln)
        text = ''.join(lines)
        if not text:
            sys.exit(f'讲稿 P{page}「{title}」没有可朗读的正文')
        sections.append(dict(page=page, title=title, text=text))
    pages = [s['page'] for s in sections]
    if pages != list(range(1, len(pages) + 1)):
        sys.exit(f'讲稿的页号不是从 1 连续到 {len(pages)}：{pages}')
    if len(sections) != expected_pages:
        sys.exit(f'讲稿 {len(sections)} 节，课件 {expected_pages} 页 —— 一页一节，对不上')
    return sections


# 让朗读更自然的替换。写讲稿时数字与公式已经用中文读法，这里只兜底几个符号。
TTS_REPLACE = [('——', '，'), ('—', '，'), ('…', '。'), ('「', ''), ('」', ''),
               ('《', ''), ('》', ''), ('·', ' '), ('≠', '不等于'),
               ('Ω', '欧米伽'), ('Θ', '希塔')]


def norm_for_tts(text):
    for a, b in TTS_REPLACE:
        text = text.replace(a, b)
    return text


# ---------------------------------------------------------------- 语音
def synth(sections, audio_dir):
    audio_dir.mkdir(parents=True, exist_ok=True)
    meta = []
    for s in sections:
        pid = '%02d' % s['page']
        spoken = norm_for_tts(s['text'])
        key = hashlib.sha256(
            ('%s|%s|%s' % (VOICE, RATE, spoken)).encode('utf-8')).hexdigest()[:16]
        mp3 = audio_dir / f'n-{pid}-{key}.mp3'
        if not mp3.exists():
            # 同一页的旧音频（讲稿改过）清掉，别在工作目录里越攒越多
            for stale in audio_dir.glob(f'n-{pid}-*.mp3'):
                stale.unlink()
            run(['edge-tts', '--voice', VOICE, '--rate', RATE,
                 '--text', spoken, '--write-media', mp3])
            mark = ' 新'
        else:
            mark = ''
        d = duration(mp3)
        meta.append(dict(page=s['page'], title=s['title'], chars=len(s['text']),
                         speech=d, clip=HEAD + d + TAIL, audio=str(mp3)))
        print('P%s %-28s %4d字  语音%6.1fs  片段%6.1fs%s'
              % (pid, s['title'][:28], len(s['text']), d, HEAD + d + TAIL, mark))
    return meta


# ---------------------------------------------------------------- 画面
def export_frames(pptx, outdir):
    """从当前 pptx 现导逐页 PNG；pptx 没变就复用。"""
    outdir.mkdir(parents=True, exist_ok=True)
    stamp = outdir / 'source.sha256'
    digest = pptx_digest(pptx)
    pngs = sorted(outdir.glob('slide-*.png'))
    if stamp.exists() and stamp.read_text().strip() == digest and pngs:
        print('画面：复用已导出的 %d 页（课件未变）' % len(pngs))
        return digest, len(pngs)
    shutil.rmtree(outdir, ignore_errors=True)
    outdir.mkdir(parents=True, exist_ok=True)
    print('画面：从当前课件导出中 …')
    run(['soffice', '--headless', '--convert-to', 'pdf', '--outdir', outdir, pptx])
    pdfs = list(outdir.glob('*.pdf'))
    if not pdfs:
        sys.exit('LibreOffice 转换失败，导不出画面')
    run(['pdftoppm', '-r', '144', '-png', pdfs[0], outdir / 'slide'])
    stamp.write_text(digest)
    n = len(list(outdir.glob('slide-*.png')))
    print('画面：导出 %d 页，来源摘要 %s' % (n, digest[:12]))
    return digest, n


# ---------------------------------------------------------------- 合成
def build_clips(meta, frames, clips):
    clips.mkdir(parents=True, exist_ok=True)
    paths = []
    for m in meta:
        pid = '%02d' % m['page']
        clip = clips / f'c-{pid}.mp4'
        run(['ffmpeg', '-y', '-loop', '1', '-framerate', '30',
             '-i', frames / f'slide-{pid}.png', '-i', m['audio'],
             '-filter_complex',
             '[0:v]scale=1920:1080:force_original_aspect_ratio=decrease,'
             f'pad=1920:1080:(ow-iw)/2:(oh-ih)/2:color={BG},setsar=1[v];'
             '[1:a]adelay=%d|%d,apad=whole_dur=%.3f,aresample=48000[a]'
             % (int(HEAD * 1000), int(HEAD * 1000), m['clip']),
             '-map', '[v]', '-map', '[a]', '-t', '%.3f' % m['clip'],
             '-c:v', 'libx264', '-preset', 'medium', '-crf', '23',
             '-pix_fmt', 'yuv420p',
             '-c:a', 'aac', '-ar', '48000', '-ac', '2', '-b:a', '128k',
             '-movflags', '+faststart', clip])
        paths.append(clip)
        print('clip', pid, 'ok')
    return paths


def concat(paths, out):
    lst = WORK / 'concat.txt'
    lst.write_text('\n'.join("file '%s'" % p for p in paths) + '\n',
                   encoding='utf-8')
    # 响度归一化并进这一步：TTS 原始输出偏轻，教室放映听不清。
    # 手工补做的步骤不可复现，重录一次就丢 —— 必须写在脚本里。
    run(['ffmpeg', '-y', '-f', 'concat', '-safe', '0', '-i', lst,
         '-c:v', 'libx264', '-profile:v', 'high', '-level:v', '4.0',
         '-pix_fmt', 'yuv420p', '-r', '30', '-fps_mode', 'cfr', '-crf', '23',
         '-video_track_timescale', '90000',
         '-af', 'loudnorm=I=-16:TP=-1.5:LRA=11',
         '-c:a', 'aac', '-ar', '48000', '-ac', '2', '-b:a', '128k',
         '-movflags', '+faststart', out])


def make_preview(src, dst):
    """从成品转一份 720p 单声道的轻量版。

    成品一章约 50MB，很多渠道（邮件、聊天工具、上传表单）传不动。
    做成脚本里的一档而不是临时敲一行 ffmpeg：临时命令不可复现，
    下次想再要一份就得凭记忆重敲，参数一定会漂。
    画质档位有意选得保守 —— 它只用来看清讲了什么，正式放映用成品。
    """
    dst.parent.mkdir(parents=True, exist_ok=True)
    run(['ffmpeg', '-y', '-i', src,
         '-vf', 'scale=1280:720:flags=lanczos',
         '-c:v', 'libx264', '-preset', 'medium', '-crf', '28',
         '-pix_fmt', 'yuv420p', '-r', '30', '-fps_mode', 'cfr',
         '-c:a', 'aac', '-ar', '44100', '-ac', '1', '-b:a', '64k',
         '-movflags', '+faststart', dst])


# ---------------------------------------------------------------- 产出
def write_srt(meta, path):
    path.parent.mkdir(parents=True, exist_ok=True)
    lines, t = [], 0.0
    for i, m in enumerate(meta, start=1):
        lines.append('%d\n%s --> %s\nP%d｜%s\n'
                     % (i, srt_ts(t + HEAD), srt_ts(t + HEAD + m['speech']),
                        m['page'], m['title']))
        t += m['clip']
    path.write_text('\n'.join(lines), encoding='utf-8')


def write_timeline(meta, path, pptx_sha, script_sha, total):
    path.parent.mkdir(parents=True, exist_ok=True)
    t, rows = 0.0, []
    for m in meta:
        rows.append(dict(page=m['page'], title=m['title'], start=fmt(t),
                         end=fmt(t + m['clip']), seconds=round(m['clip'], 1),
                         chars=m['chars']))
        t += m['clip']
    path.write_text(json.dumps(
        dict(voice=VOICE, rate=RATE, total=fmt(total),
             source_pptx_sha256=pptx_sha, source_script_sha256=script_sha,
             sections=rows), ensure_ascii=False, indent=2), encoding='utf-8')


TIMETABLE_HEAD = '## 附：时间控制表'


def narration_body(path):
    """讲稿里**手写的那一半**（时间控制表之前，且不含表前的那条分隔线）。

    ⚠️ 必须归一化到「可重复」的形状，否则 write_timetable 每写一次，
    文末的分隔线就多一条、摘要就变一次 —— 视频从此永远显示「过期」，
    而实际上一个字都没改。第一次跑就踩了这个坑。
    """
    body = io.open(path, encoding='utf-8').read().split(TIMETABLE_HEAD)[0]
    return re.sub(r'\n+-{3,}\s*$', '', body.rstrip()).rstrip()


def write_timetable(script_path, meta, total):
    """把实测时长写回讲稿文末。手写的是词，机器写的是秒。"""
    head = narration_body(script_path)
    rows = ['%s\n' % TIMETABLE_HEAD,
            '<!-- 由 make_video.py 按实测时长自动生成，不要手改 -->\n',
            '语音 `%s`，语速 `%s`，合计 **%s**（%d 页，共 %d 字）。\n'
            % (VOICE, RATE, fmt(total), len(meta), sum(m['chars'] for m in meta)),
            '| 页 | 标题 | 起 | 止 | 时长 | 字数 |',
            '| ---: | --- | ---: | ---: | ---: | ---: |']
    t = 0.0
    for m in meta:
        rows.append('| %d | %s | %s | %s | %.1fs | %d |'
                    % (m['page'], m['title'], fmt(t), fmt(t + m['clip']),
                       m['clip'], m['chars']))
        t += m['clip']
    io.open(script_path, 'w', encoding='utf-8').write(
        head + '\n\n---\n\n' + '\n'.join(rows) + '\n')


# ---------------------------------------------------------------- --check
def check(chapter, pptx, script, mp4, timeline):
    """校验时间轴记录的来源摘要是否仍等于当前的课件与讲稿。

    成品 mp4 不入库（一章约 50MB），所以它缺席**不判红**，只提示可重建 ——
    和 tools/pdfref.py 缺少扫描件时的做法一致。判红的是「时间轴与来源对不上」：
    那说明课件或讲稿改过而视频没重建，此时哪怕 mp4 在，它也是旧的。
    """
    problems, hints = [], []
    for label, path in (('课件', pptx), ('讲稿', script), ('时间轴', timeline)):
        if not path.exists():
            problems.append(f'{label}不存在：{path.name}')
    if problems:
        return problems, hints
    tl = json.loads(timeline.read_text(encoding='utf-8'))
    if tl.get('source_pptx_sha256') != pptx_digest(pptx):
        problems.append(f'第{chapter}章视频过期：课件改过而视频没重建')
    if tl.get('source_script_sha256') != script_digest(script):
        problems.append(f'第{chapter}章视频过期：讲稿改过而视频没重建')
    pages_in_script = len(SECTION.findall(
        script.read_text(encoding='utf-8').split(TIMETABLE_HEAD)[0]))
    if len(tl.get('sections', [])) != pages_in_script:
        problems.append(f'第{chapter}章时间轴 {len(tl.get("sections", []))} 页，'
                        f'讲稿 {pages_in_script} 页 —— 对不上')
    if not mp4.exists():
        hints.append(f'第{chapter}章成品视频不在本地（按设计不入库）：'
                     f'`python3 make_video.py {chapter}` 可重建')
    return problems, hints


def script_digest(path):
    """只对讲稿的**朗读部分**取摘要 —— 文末的时间控制表是这份脚本自己写回去的，
    把它算进来，视频就会永远显示「过期」。"""
    return hashlib.sha256(narration_body(path).encode('utf-8')).hexdigest()


# ---------------------------------------------------------------- 入口
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('chapter', help='章号，如 01')
    ap.add_argument('--audio-only', action='store_true', help='只跑 TTS 看总时长')
    ap.add_argument('--check', action='store_true', help='只校验产物是否最新')
    ap.add_argument('--preview', action='store_true',
                    help='从成品转一份 720p 轻量版，便于传阅')
    args = ap.parse_args()

    ch = args.chapter.zfill(2)
    if ch not in build_all.CHAPTERS:
        sys.exit(f'未知章号：{args.chapter}')
    stem = build_all.CHAPTERS[ch]
    pptx = HERE / (stem + '.pptx')
    script = HERE / 'content' / f'ch{ch}_narration.md'
    mp4 = HERE / (stem + '.mp4')
    srt = HERE / 'video' / f'ch{ch}.srt'
    timeline = HERE / 'video' / f'ch{ch}.timeline.json'

    if args.preview:
        if not mp4.exists():
            sys.exit(f'成品不存在：{mp4.name}。先跑 python3 make_video.py {ch}')
        preview = HERE / 'video' / f'ch{ch}-preview.mp4'
        make_preview(mp4, preview)
        print('输出：%s　%.1f MB（成品 %.1f MB）'
              % (preview.name, preview.stat().st_size / 1e6,
                 mp4.stat().st_size / 1e6))
        return 0

    if args.check:
        problems, hints = check(ch, pptx, script, mp4, timeline)
        for p in problems:
            print('❌', p)
        for h in hints:
            print('·', h)
        if not problems:
            tl = json.loads(timeline.read_text(encoding='utf-8'))
            print('✅ 第%s章视频与课件、讲稿同版：%d 页，总时长 %s'
                  % (ch, len(tl['sections']), tl['total']))
        return 1 if problems else 0

    if not pptx.exists():
        sys.exit(f'课件不存在：{pptx}。先跑 python3 build_all.py {ch}')
    if not script.exists():
        sys.exit(f'讲稿不存在：{script}')

    from pptx import Presentation
    pages = len(Presentation(str(pptx)).slides)
    sections = parse_narration(script, pages)

    global WORK
    WORK = WORK / f'ch{ch}'
    meta = synth(sections, WORK / 'audio')
    total = sum(m['clip'] for m in meta)
    print('\n合计 %s（%.1f 秒），%d 页，共 %d 字，voice=%s rate=%s'
          % (fmt(total), total, len(meta), sum(m['chars'] for m in meta),
             VOICE, RATE))
    if args.audio_only:
        return 0

    pptx_sha, nframes = export_frames(pptx, WORK / 'frames')
    if nframes != pages:
        sys.exit(f'导出 {nframes} 张画面，课件 {pages} 页 —— 对不上')

    build_clips(meta, WORK / 'frames', WORK / 'clips')
    concat([WORK / 'clips' / ('c-%02d.mp4' % m['page']) for m in meta], mp4)

    write_srt(meta, srt)
    write_timeline(meta, timeline, pptx_sha, script_digest(script), total)
    write_timetable(script, meta, total)
    size = mp4.stat().st_size / 1e6
    print('\n输出：%s　总时长 %s　%.1f MB' % (mp4.name, fmt(total), size))
    return 0


if __name__ == '__main__':
    sys.exit(main())
