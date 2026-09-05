#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""verify.py —— 讲义与课件的一致性闸门。

用法:
    python3 verify.py             # 第 1–6 项，几秒
    python3 verify.py --render    # 加第 7 项渲染检查（需 libreoffice + poppler-utils）

检查项:
  1 配对    每章 content/chNN.py + 同名 .md + 同名 .pptx，三者齐全
  2 元数据  讲义有一级标题 / Updated / Compiled by / 仓库 URL；META 四个键齐全
  3 资源    课件 image 页引用的图、讲义里的本地相对链接，都真的存在
  4 代码    content/ 与 deck.py 能 ast.parse；courseware/code/ 下每个 .cpp
            在 -Werror 下真编译、真运行，退出码为 0
  5 可重生成  课件能从 content/ 重建；页数与 README 的表格一致，
            且重建产物与已提交的 .pptx **逐段文本相同**（防止改了源没重建）
  6 版面标记  放映稿的非等宽文字里不得印出 `**` 或反引号；条目不得以空白开头
  7 渲染    逐页检查文字未越出版心、未侵入页脚、**没有两段文字压在一起**；
            PDF 已嵌入中文字体（--render）

⚠️ 第 6 项的存在理由：`**强调**` 与 `` `等宽` `` 是给排版引擎看的标记，不是内容。
   漏掉一层解析，它们就原样印在投影幕上；人眼对一个反引号极不敏感，
   所以这件事必须由代码保证 —— cs101 曾经三轮逐页人工复核全部漏过。

只用标准库 + python-pptx。退出码 0 表示全绿。
"""

import argparse
import ast
import re
import subprocess
import sys
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))
sys.path.insert(0, str(HERE / 'content'))

import build_all                                    # noqa: E402
import deck                                         # noqa: E402
from pptx import Presentation                       # noqa: E402
from pptx.oxml.ns import qn                         # noqa: E402

FAILURES = []
NOTES = []


def fail(item, msg):
    FAILURES.append(f'❌ [{item}] {msg}')


def note(msg):
    NOTES.append(f'· {msg}')


def written_chapters():
    """已经写了内容的章（content/chNN.py 存在的那些）。"""
    return [ch for ch in sorted(build_all.CHAPTERS)
            if (HERE / 'content' / f'ch{ch}.py').exists()]


# ---------------------------------------------------------------- 1 配对
def check_pairing(chapters):
    if not chapters:
        fail('1 配对', 'content/ 下一个 chNN.py 都没有')
        return
    for ch in chapters:
        stem = build_all.CHAPTERS[ch]
        for suffix in ('.md', '.pptx'):
            if not (HERE / (stem + suffix)).is_file():
                fail('1 配对', f'第{ch}章缺 {stem}{suffix}')
    note(f'已编写 {len(chapters)} 章：' + '、'.join('第%s章' % c.lstrip('0')
                                                   for c in chapters))


# ---------------------------------------------------------------- 2 元数据
REPO_URL = 'https://github.com/GMyhf/dsa-modernization'
META_KEYS = ('title', 'subtitle', 'footer', 'info')


def check_metadata(chapters, modules):
    for ch in chapters:
        stem = build_all.CHAPTERS[ch]
        md = (HERE / (stem + '.md')).read_text(encoding='utf-8')
        head = md[:900]
        for label, pat in (('一级标题', r'^# .+'),
                           ('Updated 时间戳', r'\*Updated \d{4}-\d{2}-\d{2}'),
                           ('Compiled by', r'\*Compiled by '),
                           ('仓库 URL', re.escape(REPO_URL))):
            if not re.search(pat, head, re.M):
                fail('2 元数据', f'{stem}.md 开头缺{label}')
        meta = modules[ch].META
        missing = [k for k in META_KEYS if not meta.get(k)]
        if missing:
            fail('2 元数据', f'content/ch{ch}.py 的 META 缺 {missing}')
        # 讲义与课件必须指向同一章
        chapter_no = str(int(ch))
        if f'第{chapter_no}章' not in meta.get('title', ''):
            fail('2 元数据', f"content/ch{ch}.py 的 META['title'] 里没有「第{chapter_no}章」")
        if not md.lstrip().startswith(f'# 第{chapter_no}章'):
            fail('2 元数据', f'{stem}.md 的一级标题不是「第{chapter_no}章 …」')


# ---------------------------------------------------------------- 3 资源
LINK = re.compile(r'\[[^\]]*\]\(([^)\s]+)\)')
IMG = re.compile(r'!\[[^\]]*\]\(([^)\s]+)\)')


def check_assets(chapters, modules):
    for ch in chapters:
        stem = build_all.CHAPTERS[ch]
        # 课件 image 页
        for spec in modules[ch].SLIDES:
            if spec[0] == 'image':
                p = HERE / spec[2]
                if not p.is_file():
                    fail('3 资源', f'content/ch{ch}.py 的 image 页「{spec[1]}」'
                                   f'找不到图片：{spec[2]}')
        # 讲义里的本地相对链接
        md_path = HERE / (stem + '.md')
        for m in set(LINK.findall(md_path.read_text(encoding='utf-8'))
                     + IMG.findall(md_path.read_text(encoding='utf-8'))):
            if m.startswith(('http://', 'https://', '#', 'mailto:')):
                continue
            target = (md_path.parent / m.split('#', 1)[0]).resolve()
            if not target.exists():
                fail('3 资源', f'{stem}.md 的链接指向不存在的路径：{m}')


# ---------------------------------------------------------------- 4 代码
def check_code(chapters):
    for py in sorted(HERE.glob('*.py')) + sorted((HERE / 'content').glob('*.py')):
        try:
            ast.parse(py.read_text(encoding='utf-8'))
        except SyntaxError as exc:
            fail('4 代码', f'{py.relative_to(HERE)} 语法错误：{exc}')

    sources = sorted((HERE / 'code').rglob('*.cpp')) if (HERE / 'code').is_dir() else []
    if not sources:
        note('courseware/code/ 下没有 .cpp，跳过编译')
        return
    with tempfile.TemporaryDirectory() as tmp:
        for src in sources:
            exe = Path(tmp) / src.stem
            build = subprocess.run(
                ['c++', '-std=c++17', '-Wall', '-Wextra', '-Werror',
                 str(src), '-o', str(exe)],
                capture_output=True, text=True)
            if build.returncode != 0:
                fail('4 代码', f'{src.relative_to(HERE)} 编译失败：\n'
                               + build.stderr.strip()[:1200])
                continue
            run = subprocess.run([str(exe)], capture_output=True, text=True,
                                 timeout=60)
            if run.returncode != 0:
                fail('4 代码', f'{src.relative_to(HERE)} 运行失败，'
                               f'退出码 {run.returncode}：\n'
                               + (run.stderr or run.stdout).strip()[:800])
    note(f'编译并运行 {len(sources)} 个 .cpp（-Werror），全部退出码 0')


# ---------------------------------------------------------------- 5 可重生成
README_ROW = re.compile(r'^\|\s*(\d+)\s*\|\s*`([^`]+)`\s*\|\s*(\d+)\s*\|', re.M)


def slide_texts(path):
    """把一份 pptx 摊平成「每页的文字列表」，用于逐段比对。"""
    out = []
    for slide in Presentation(str(path)).slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            elif shape.has_table:
                texts += [c.text for r in shape.table.rows for c in r.cells]
        out.append(texts)
    return out


def check_regenerate(chapters, modules):
    readme = HERE / 'README.md'
    declared = {}
    if readme.is_file():
        for num, stem, pages in README_ROW.findall(
                readme.read_text(encoding='utf-8')):
            declared[f'{int(num):02d}'] = (stem, int(pages))
    else:
        fail('5 可重生成', 'README.md 不存在')

    with tempfile.TemporaryDirectory() as tmp:
        for ch in chapters:
            stem = build_all.CHAPTERS[ch]
            rebuilt = Path(tmp) / (stem + '.pptx')
            pages = deck.build(modules[ch].META, modules[ch].SLIDES, str(rebuilt))
            committed = HERE / (stem + '.pptx')
            if not committed.is_file():
                continue
            if slide_texts(rebuilt) != slide_texts(committed):
                fail('5 可重生成',
                     f'{stem}.pptx 与 content/ch{ch}.py 不一致 —— '
                     f'改了源没重建？跑 `python3 build_all.py {ch}`')
            if ch not in declared:
                fail('5 可重生成', f'README.md 的文件清单里没有第{ch}章')
            else:
                dstem, dpages = declared[ch]
                if dstem != stem:
                    fail('5 可重生成',
                         f'README.md 第{ch}章写的文件名是 {dstem}，实际是 {stem}')
                if dpages != pages:
                    fail('5 可重生成',
                         f'README.md 写第{ch}章 {dpages} 页，实际生成 {pages} 页')
            note(f'第{ch}章 {pages} 页，与 content/ch{ch}.py 逐段一致')


# ---------------------------------------------------------------- 6 版面标记
LEAK = re.compile(r'\*\*|`')
LEADING_BLANK = re.compile(r'^[\s　]')


def check_markup(chapters, modules):
    for ch in chapters:
        stem = build_all.CHAPTERS[ch]
        path = HERE / (stem + '.pptx')
        if not path.is_file():
            continue
        for pno, slide in enumerate(Presentation(str(path)).slides, start=1):
            for shape in slide.shapes:
                frames = []
                if shape.has_text_frame:
                    frames.append(shape.text_frame)
                elif shape.has_table:
                    frames += [c.text_frame for r in shape.table.rows
                               for c in r.cells]
                for tf in frames:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.name == deck.MONO_FONT:
                                continue    # 代码 / 示意图原样输出，不解析标记
                            if LEAK.search(run.text):
                                fail('6 版面标记',
                                     f'{stem}.pptx 第 {pno} 页印出了排版标记：'
                                     f'{run.text!r}')
        # 条目不得以空白（含全角空格）开头 —— 放映时会莫名其妙缩进一格
        for spec in modules[ch].SLIDES:
            items = []
            if spec[0] == 'bullets':
                items = spec[2]
            elif spec[0] == 'two':
                items = list(spec[3]) + list(spec[5])
            for raw in items:
                body = raw[2:] if raw.startswith('- ') else raw
                if LEADING_BLANK.match(body):
                    fail('6 版面标记',
                         f'content/ch{ch}.py 的条目以空白开头：{raw!r}')


# ---------------------------------------------------------------- 7 渲染
BODY_BOTTOM_PT = (1.42 + 5.32) * 72          # 版心底 = 485.28pt（540pt 版面）
FOOTER_TEXT_TOP_PT = 496.0                   # 页脚文本实测顶端，比文本框理论值略高
BODY_TOL_PT = 4.0                            # 字形下伸部的容差

CJK_FONT_HINT = ('CJK', 'YaHei', 'SimSun', 'SimHei', 'Song', 'Hei', 'Ming',
                 'Noto Sans SC', 'Source Han')


# 两个字框重叠到什么程度算「压字」。同一行相邻的字会轻微相接，
# 所以按「交集面积 / 较小者面积」判，而不是「是否相交」。
OVERLAP_RATIO = 0.45


def overlapping_words(words):
    """找出彼此压在一起的字框。words 是 (xMin, yMin, xMax, yMax, 文本) 列表。

    ⚠️ 这条判据是补上来的。第 1 章「六种算法设计方法」那一页，
    表格的图注**压在最后一行的字上** —— 文字全在版心内、也没碰页脚，
    越界检查一声不吭。版面事故不只有「越出去」一种，还有「叠起来」。
    """
    hits = []
    ordered = sorted(words, key=lambda w: w[1])
    for i, a in enumerate(ordered):
        for b in ordered[i + 1:]:
            if b[1] >= a[3]:            # 按 yMin 排过序，后面的都在 a 下方了
                break
            ix = min(a[2], b[2]) - max(a[0], b[0])
            iy = min(a[3], b[3]) - max(a[1], b[1])
            if ix <= 0 or iy <= 0:
                continue
            area_a = (a[2] - a[0]) * (a[3] - a[1])
            area_b = (b[2] - b[0]) * (b[3] - b[1])
            smaller = min(area_a, area_b)
            if smaller > 0 and ix * iy / smaller > OVERLAP_RATIO:
                hits.append((a[4], b[4]))
    return hits


def intrudes_into_footer(y_min_pt, y_max_pt):
    """正文文字是否越出版心底部并侵入页脚区（坐标已换算到 540pt 版面）。

    抽成具名函数，是为了能脱离 LibreOffice 直接回归这条判据 —— 它的阈值来自实测：
    页脚文本自己的 yMin ≈ 497.5，而文本框理论值是 498.2；按理论值取，
    **每一页的页脚都会被误判成溢出**。
    """
    return y_min_pt < FOOTER_TEXT_TOP_PT and y_max_pt > BODY_BOTTOM_PT + BODY_TOL_PT


def check_render(chapters):
    soffice = None
    for cand in ('soffice', 'libreoffice'):
        try:
            subprocess.run([cand, '--version'], capture_output=True, check=True)
            soffice = cand
            break
        except (FileNotFoundError, subprocess.CalledProcessError):
            continue
    if soffice is None:
        fail('7 渲染', '找不到 libreoffice / soffice')
        return
    try:
        subprocess.run(['pdftotext', '-v'], capture_output=True)
    except FileNotFoundError:
        fail('7 渲染', '找不到 pdftotext（poppler-utils）')
        return

    pptxs = [HERE / (build_all.CHAPTERS[ch] + '.pptx') for ch in chapters]
    pptxs = [p for p in pptxs if p.is_file()]
    with tempfile.TemporaryDirectory() as tmp:
        proc = subprocess.run(
            [soffice, '--headless', '--convert-to', 'pdf', '--outdir', tmp]
            + [str(p) for p in pptxs],
            capture_output=True, text=True, timeout=1800)
        pdfs = sorted(Path(tmp).glob('*.pdf'))
        if len(pdfs) != len(pptxs):
            fail('7 渲染', f'转换未全部成功：期望 {len(pptxs)} 份 PDF，'
                           f'实得 {len(pdfs)} 份；soffice 退出码 {proc.returncode}')
            return
        total = 0
        for pdf in pdfs:
            fonts = subprocess.run(['pdffonts', str(pdf)],
                                   capture_output=True, text=True).stdout
            if not any(h in fonts for h in CJK_FONT_HINT):
                fail('7 渲染', f'{pdf.name}: PDF 未嵌入任何已知中文字体，'
                               f'渲染结果可能是方框（本机字体环境限制）')
            bbox = subprocess.run(['pdftotext', '-bbox', str(pdf), '-'],
                                  capture_output=True, text=True).stdout
            pages = re.findall(
                r'<page width="([\d.]+)" height="([\d.]+)">(.*?)</page>',
                bbox, re.S)
            total += len(pages)
            for pno, (pw, ph, body) in enumerate(pages, start=1):
                pw, ph = float(pw), float(ph)
                sy = 540.0 / ph
                words = []
                for x, y, x2, y2, text in re.findall(
                        r'<word xMin="([\d.]+)" yMin="([\d.]+)" '
                        r'xMax="([\d.]+)" yMax="([\d.]+)">([^<]*)', body):
                    x, y, x2, y2 = float(x), float(y), float(x2), float(y2)
                    words.append((x, y, x2, y2, text))
                    if x2 > pw + 0.5 or y2 > ph + 0.5 or x < -0.5 or y < -0.5:
                        fail('7 渲染', f'{pdf.name} 第 {pno} 页：文字越出页面')
                        break
                    if intrudes_into_footer(y * sy, y2 * sy):
                        fail('7 渲染',
                             f'{pdf.name} 第 {pno} 页：正文越出版心底部并侵入页脚区'
                             f'（文字底 {y2 * sy:.1f}pt > 版心底 {BODY_BOTTOM_PT:.1f}pt）')
                        break
                hits = overlapping_words(words)
                if hits:
                    a, b = hits[0]
                    fail('7 渲染',
                         f'{pdf.name} 第 {pno} 页：两段文字压在一起'
                         f'（{a!r} 与 {b!r}，共 {len(hits)} 处）')
        note(f'渲染检查：{len(pdfs)} 份 PDF，共 {total} 页')


# ---------------------------------------------------------------- 入口
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--render', action='store_true', help='加做第 7 项渲染检查')
    args = ap.parse_args()

    chapters = written_chapters()
    modules = {}
    import importlib
    for ch in chapters:
        modules[ch] = importlib.import_module(f'ch{ch}')

    check_pairing(chapters)
    check_metadata(chapters, modules)
    check_assets(chapters, modules)
    check_code(chapters)
    check_regenerate(chapters, modules)
    check_markup(chapters, modules)
    if args.render:
        check_render(chapters)

    for line in NOTES:
        print(line)
    if FAILURES:
        print()
        for line in FAILURES:
            print(line)
        print(f'\n{len(FAILURES)} 项不通过。')
        return 1
    print('\n全绿。')
    return 0


if __name__ == '__main__':
    sys.exit(main())
